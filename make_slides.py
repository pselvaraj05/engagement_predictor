"""Generate presentation slide deck using the EveryoneSocial template."""
import os
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = "/Users/prashanthselvaraj/Github/engagement_predictor/Engagement prediction.pptx"
fig_dir = "./figures"

# ─── Template Colors ─────────────────────────────────────────────────────────
BLUE = RGBColor(0x38, 0x9A, 0xFF)
YELLOW = RGBColor(0xF9, 0xA8, 0x1A)
RED = RGBColor(0xEA, 0x1D, 0x64)
DARK = RGBColor(0x3A, 0x38, 0x38)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_HEADER_BG = RGBColor(0x38, 0x9A, 0xFF)
TABLE_ROW_BG = RGBColor(0xF0, 0xF4, 0xF8)

# ─── Load Template, strip slides, save as clean base, then reload ─────────────
import tempfile, shutil, zipfile, re

def make_clean_template(src_path):
    """Create a copy of the template with all slides removed at the ZIP level."""
    tmp = tempfile.mktemp(suffix=".pptx")
    with zipfile.ZipFile(src_path, "r") as zin:
        # Read files that need patching first
        all_files = {}
        skip_pattern = re.compile(
            r"^ppt/slides/|^ppt/notesSlides/|^ppt/slides/_rels/"
        )
        patch_names = {"ppt/presentation.xml", "[Content_Types].xml", "ppt/_rels/presentation.xml.rels"}

        for item in zin.infolist():
            if skip_pattern.match(item.filename):
                continue
            all_files[item.filename] = zin.read(item.filename)

        # Patch presentation.xml
        pres_xml = all_files["ppt/presentation.xml"].decode("utf-8")
        pres_xml = re.sub(r'<p:sldIdLst>.*?</p:sldIdLst>', '<p:sldIdLst/>', pres_xml, flags=re.DOTALL)
        all_files["ppt/presentation.xml"] = pres_xml.encode("utf-8")

        # Patch [Content_Types].xml
        ct_xml = all_files["[Content_Types].xml"].decode("utf-8")
        ct_xml = re.sub(r'<Override[^>]*PartName="/ppt/slides/slide\d+\.xml"[^>]*/>', '', ct_xml)
        ct_xml = re.sub(r'<Override[^>]*PartName="/ppt/notesSlides/[^"]*"[^>]*/>', '', ct_xml)
        all_files["[Content_Types].xml"] = ct_xml.encode("utf-8")

        # Patch ppt/_rels/presentation.xml.rels
        rels_xml = all_files["ppt/_rels/presentation.xml.rels"].decode("utf-8")
        rels_xml = re.sub(r'<Relationship[^>]*Target="slides/slide\d+\.xml"[^>]*/>', '', rels_xml)
        all_files["ppt/_rels/presentation.xml.rels"] = rels_xml.encode("utf-8")

    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for fname, data in all_files.items():
            zout.writestr(fname, data)

    return tmp

clean_template = make_clean_template(TEMPLATE)
prs = Presentation(clean_template)
os.remove(clean_template)

# Layout references
LY_BLANK = prs.slide_layouts[11]         # BLANK


# ─── Helpers ──────────────────────────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text, size=15,
                color=DARK, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=15,
                color=DARK, bold=False, spacing=Pt(6)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = spacing
    return tb


def add_card(slide, left, top, width, height, fill=CARD_BG, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_line(slide, left, top, width, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.65),
                title, size=30, color=DARK, bold=True)
    add_line(slide, Inches(0.7), Inches(0.95), Inches(2))
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11), Inches(0.45),
                    subtitle, size=14, color=GRAY)


def metric_card(slide, left, top, width, height, label, value, accent=BLUE):
    add_card(slide, left, top, width, height, fill=WHITE, border=accent)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
                label, size=11, color=GRAY)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.55),
                value, size=26, color=accent, bold=True)


def section_label(slide, text, left=Inches(0.7), top=Inches(6.6)):
    """Small section label at bottom-left, matching template style."""
    add_textbox(slide, left, top, Inches(3), Inches(0.3), text, size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(1.2),
            "Engagement Prediction", size=44, color=DARK, bold=True)
add_line(slide, Inches(0.8), Inches(2.8), Inches(2.5), BLUE)
add_textbox(slide, Inches(0.8), Inches(3.1), Inches(6), Inches(0.8),
            "Predicting who engages and how audiences\ndistribute across shares",
            size=20, color=GRAY)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
            "Methods & Results", size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Problem Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Problem Overview")
section_label(slide, "Overview")

cards = [
    ("Model 1: Engagement Volume", BLUE,
     ["How many engagements will a share receive?",
      "Hierarchical binary + tier classification",
      "Tier-specific count regression",
      "Output: Low / Medium / High / Viral + count"]),
    ("Model 2: Engager Profile", RED,
     ["Who will engage with this share?",
      "Multi-class classification per target",
      "Predicts job title, industry, company",
      "Output: Top-K ranked predictions per share"]),
    ("Model 3: Audience Distribution", YELLOW,
     ["What audience mix will a share attract?",
      "Multi-output regression on proportions",
      "Predicts % breakdown across categories",
      "Output: Probability distribution per target"]),
]

for i, (title, accent, bullets) in enumerate(cards):
    left = Inches(0.5 + i * 4.15)
    add_card(slide, left, Inches(1.6), Inches(3.9), Inches(4.8), fill=CARD_BG, border=accent)
    add_line(slide, left + Inches(0.25), Inches(1.85), Inches(1.5), accent)
    add_textbox(slide, left + Inches(0.25), Inches(2.0), Inches(3.4), Inches(0.5),
                title, size=16, color=accent, bold=True)
    add_bullets(slide, left + Inches(0.25), Inches(2.7), Inches(3.4), Inches(3.3),
                bullets, size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Data Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Data Overview", "LinkedIn engagement data pipeline")
section_label(slide, "Data")

sources = [
    ("Engagements", "52.4M rows", "user_id, share_id,\nengagement_type,\nprofile_id, created_at"),
    ("Shares", "4.8M rows", "share_id, shared_at,\ncontent_type,\nuser_commentary"),
    ("Profiles", "9.0M rows", "profile_id, industry,\njob_title_role,\nemployer, location"),
    ("Clients", "1,516 rows", "client_id, title,\nstatus"),
]

for i, (name, count, desc) in enumerate(sources):
    left = Inches(0.5 + i * 3.15)
    add_card(slide, left, Inches(1.6), Inches(2.95), Inches(2.3))
    add_textbox(slide, left + Inches(0.2), Inches(1.75), Inches(2.5), Inches(0.35),
                name, size=18, color=BLUE, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(2.1), Inches(2.5), Inches(0.3),
                count, size=14, color=RED, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(2.5), Inches(2.5), Inches(1.0),
                desc, size=11, color=GRAY)

add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.2))
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5), Inches(0.35),
            "Key Data Challenge", size=18, color=RED, bold=True)
add_bullets(slide, Inches(0.8), Inches(4.85), Inches(11.5), Inches(1.5), [
    "Only 7.5% of engagements link to profiles with all 3 targets (job title + industry + company) known",
    "52.4M engagements --> 3.9M usable rows after filtering to complete profiles",
    "Rare label collapsing: categories below min frequency threshold merged into 'other'",
    "Time-based train/test split (not random) to prevent data leakage",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Feature Engineering
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Feature Engineering", "Shared across all models")
section_label(slide, "Features")

feat_groups = [
    ("Text Features", BLUE,
     ["TF-IDF on user_commentary", "Up to 2,000 features", "Unigrams + bigrams", "English stop words removed"]),
    ("Content Signals", RED,
     ["Text length, word count", "Has question mark?", "Has URL?", "Exclamation / hashtag / @"]),
    ("Temporal Features", YELLOW,
     ["Hour of day shared", "Day of week", "Is weekend?", "Is business hours (9-5)?"]),
    ("Sharer Profile", RGBColor(0x2E, 0xCC, 0x71),
     ["Industry, job title role", "Job title class & level", "Location (region, country)", "Share content type"]),
]

for i, (title, accent, bullets) in enumerate(feat_groups):
    left = Inches(0.5 + i * 3.15)
    add_card(slide, left, Inches(1.6), Inches(2.95), Inches(3.5))
    add_line(slide, left + Inches(0.2), Inches(1.8), Inches(1), accent)
    add_textbox(slide, left + Inches(0.2), Inches(1.95), Inches(2.5), Inches(0.35),
                title, size=16, color=accent, bold=True)
    add_bullets(slide, left + Inches(0.2), Inches(2.45), Inches(2.5), Inches(2.2),
                bullets, size=13, color=DARK)

add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8))
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.6),
            "All features encoded via OneHotEncoder (categorical) + StandardScaler (numeric) + TF-IDF (text), "
            "combined into a single sparse matrix. LightGBM classifiers/regressors used throughout.",
            size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Model 1 — Method
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 1: Hierarchical Engagement Prediction",
            "Predicting how many engagements a share will receive")
section_label(slide, "Model 1 - Method")

# Three stage cards
stages = [
    ("Stage 1: Binary", BLUE,
     ["Popular vs Not Popular (>50)", "LightGBM, 500 trees", "scale_pos_weight=2.0"]),
    ("Stage 2: Sub-Tier", RED,
     ["Not Popular: Low vs Medium", "Popular: High vs Viral", "LightGBM, 300 trees each"]),
    ("Stage 3: Count", YELLOW,
     ["Tier-specific regressors", "Tweedie objective", "Clipped to tier bounds"]),
]

for i, (title, accent, bullets) in enumerate(stages):
    left = Inches(0.5 + i * 4.15)
    add_card(slide, left, Inches(1.6), Inches(3.8), Inches(1.8), fill=CARD_BG, border=accent)
    add_textbox(slide, left + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.35),
                title, size=16, color=accent, bold=True)
    add_bullets(slide, left + Inches(0.2), Inches(2.15), Inches(3.4), Inches(1.0),
                bullets, size=13, color=DARK)

# Arrows
for x in [Inches(4.45), Inches(8.6)]:
    add_textbox(slide, x, Inches(2.2), Inches(0.5), Inches(0.4),
                "-->", size=22, color=LIGHT_GRAY, bold=True)

# Tier bar
add_card(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.9))
add_textbox(slide, Inches(0.7), Inches(3.9), Inches(11.5), Inches(0.3),
            "Engagement Tiers", size=14, color=DARK, bold=True)
add_textbox(slide, Inches(0.7), Inches(4.2), Inches(11.5), Inches(0.35),
            "Low (0-5)    |    Medium (6-50)    |    High (51-500)    |    Viral (500+)",
            size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Training config
add_card(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.4))
add_textbox(slide, Inches(0.7), Inches(5.1), Inches(5), Inches(0.3),
            "Training Configuration", size=14, color=DARK, bold=True)
add_bullets(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(0.9), [
    "1M stratified sample (20% Low, 30% Med, 30% High, 20% Viral)",
    "Time-based split: 524K train / 131K test",
    "Oversampled rare tiers for better representation",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Model 1 — Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 1: Results", "Hierarchical Engagement Prediction")
section_label(slide, "Model 1 - Results")

metric_card(slide, Inches(0.3), Inches(1.5), Inches(2.4), Inches(1.0),
            "Popularity ROC-AUC", "0.950", BLUE)
metric_card(slide, Inches(2.85), Inches(1.5), Inches(2.4), Inches(1.0),
            "Popular Recall", "95.3%", RGBColor(0x2E, 0xCC, 0x71))
metric_card(slide, Inches(5.4), Inches(1.5), Inches(2.4), Inches(1.0),
            "Popular F1", "0.724", YELLOW)
metric_card(slide, Inches(7.95), Inches(1.5), Inches(2.4), Inches(1.0),
            "Tier Accuracy", "72.5%", BLUE)
metric_card(slide, Inches(10.5), Inches(1.5), Inches(2.4), Inches(1.0),
            "Count MAE", "49.8", RED)

add_card(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(3.7))
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.35),
            "Key Findings", size=18, color=DARK, bold=True)
add_bullets(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.0), [
    "ROC-AUC of 0.950 -- strong discrimination between popular and non-popular posts",
    "95.3% recall on popular posts -- catches nearly all high-performing content",
    "Precision tradeoff at 58.4% -- some false positives, acceptable for recall-first strategy",
    "72.5% tier accuracy across 4 tiers -- much better than random (25%) or majority-class baseline",
    "Count MAE of 49.8 -- reasonable given engagement counts range from 0 to 50K+",
    "Viral tier is hardest to predict precisely due to extreme outlier behavior",
], size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 7-10: Model 1 — Individual Figure Slides
# ═══════════════════════════════════════════════════════════════════════════════
slide_figs_dir = "./figures/slides"
model1_figs = [
    ("roc_curve.png", "ROC Curve: Popularity Detection", "Model 1 - ROC"),
    ("confusion_matrix_counts.png", "4-Tier Confusion Matrix", "Model 1 - Confusion Matrix"),
    ("confusion_matrix_recall.png", "4-Tier Confusion Matrix (Recall %)", "Model 1 - Recall"),
    ("tier_distribution.png", "Actual vs Predicted Tier Distribution", "Model 1 - Distribution"),
]

for fname, title, label in model1_figs:
    fpath = os.path.join(slide_figs_dir, fname)
    if os.path.exists(fpath):
        slide = prs.slides.add_slide(LY_BLANK)
        slide_title(slide, title)
        section_label(slide, label)
        slide.shapes.add_picture(fpath, Inches(1.5), Inches(1.3), Inches(10), Inches(5.8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Model 2 — Method
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 2: Engager Profile Prediction",
            "Predicting who will engage with a share")
section_label(slide, "Model 2 - Method")

# Left: Approach
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.35),
            "Approach", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.3), Inches(4.2), [
    "Training data at engagement level: one row per (share, engager) pair",
    "Three independent LightGBM classifiers, one per target:",
    "     Job Title Role (13 classes after collapsing rare)",
    "     Industry (16 classes)",
    "     Company (47 classes via employer_client_id)",
    "",
    "Pre-filter profiles to those with ALL 3 targets known",
    "Rare labels below threshold collapsed to 'other'",
    "     Job title/Industry: min 500 occurrences",
    "     Company: min 200 occurrences",
    "",
    "LightGBM: 500 trees, lr=0.05, 127 leaves",
    "Time-based 80/20 train/test split",
], size=13, color=DARK)

# Right: Metrics + Pipeline
add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.7))
add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.35),
            "Top-K Accuracy Metric", size=18, color=RED, bold=True)
add_bullets(slide, Inches(6.9), Inches(2.05), Inches(5.7), Inches(1.8), [
    "Top-1: Was the most confident prediction correct?",
    "Top-3: Was true label in the 3 most confident?",
    "Top-5: Was true label in the 5 most confident?",
    "",
    "Knowing top 3 likely segments is as actionable as the exact one.",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(4.4), Inches(6.2), Inches(2.2))
add_textbox(slide, Inches(6.9), Inches(4.5), Inches(5.5), Inches(0.35),
            "Data Pipeline", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(6.9), Inches(4.95), Inches(5.7), Inches(1.3), [
    "52.4M total engagements",
    "--> 3.9M with complete engager profile (7.5%)",
    "--> 50K sampled for training (scalable to full)",
    "--> 32K train / 8K test after time split",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Model 2 — Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 2: Results", "Engager Profile Prediction (50K sample)")
section_label(slide, "Model 2 - Results")

# Results table using shapes
targets = [
    ("Job Role", "13", "0.304", "0.546", "0.710", RED),
    ("Industry", "16", "0.417", "0.726", "0.868", BLUE),
    ("Company", "47", "0.829", "0.908", "0.931", RGBColor(0x2E, 0xCC, 0x71)),
]

# Header row
headers = ["Target", "Classes", "Top-1", "Top-3", "Top-5"]
col_lefts = [Inches(0.8), Inches(3.5), Inches(5.3), Inches(7.3), Inches(9.3)]
table_top = Inches(1.5)

add_card(slide, Inches(0.5), table_top, Inches(12.3), Inches(0.55), fill=BLUE)
for pos, hdr in zip(col_lefts, headers):
    add_textbox(slide, pos, table_top + Inches(0.08), Inches(2), Inches(0.35),
                hdr, size=15, color=WHITE, bold=True)

# Data rows
for i, (name, classes, t1, t3, t5, accent) in enumerate(targets):
    row_top = table_top + Inches(0.6 + i * 0.6)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, Inches(0.5), row_top, Inches(12.3), Inches(0.55), fill=bg)
    vals = [name, classes, t1, t3, t5]
    colors = [accent, DARK, DARK, BLUE, RGBColor(0x2E, 0xCC, 0x71)]
    bolds = [True, False, False, True, True]
    sizes = [15, 15, 17, 17, 17]
    for pos, val, c, b, sz in zip(col_lefts, vals, colors, bolds, sizes):
        add_textbox(slide, pos, row_top + Inches(0.08), Inches(2), Inches(0.35),
                    val, size=sz, color=c, bold=b)

# Insights
add_card(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(3.0))
add_textbox(slide, Inches(0.8), Inches(3.7), Inches(11), Inches(0.35),
            "Analysis", size=18, color=DARK, bold=True)
add_bullets(slide, Inches(0.8), Inches(4.15), Inches(11.5), Inches(2.2), [
    "Company prediction is strongest (83% top-1, 93% top-5) -- employer_client_id provides strong signal",
    "Industry prediction solid at top-3 (73%) -- share content and sharer profile correlate with engager industry",
    "Job role is hardest (30% top-1) but reaches 71% at top-5 -- roles are inherently more diverse",
    "All targets benefit significantly from top-K: knowing the top 3-5 likely segments is very actionable",
    "Results on 50K sample -- full dataset run expected to improve further",
], size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Power Analysis — Method
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Power Analysis: Sample Size Requirements",
            "If we stop receiving LinkedIn data and scrape instead, how much do we need?")
section_label(slide, "Power Analysis - Method")

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.35),
            "Simulation Design", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.3), Inches(4.2), [
    "Use existing dataset as ground truth",
    "Simulate having fewer shares and engagements:",
    "     Part A: Vary number of shares (50 to 20,000)",
    "     Part B: Vary engagements per share (1 to 10)",
    "     Part C: Grid search (shares x engagements)",
    "",
    "For each configuration:",
    "     Sample shares and engagements from real data",
    "     Train LightGBM classifiers on the sample",
    "     Evaluate on held-out test set",
    "     Repeat 3 trials for confidence intervals",
    "",
    "Fixed label sets computed once from full data",
    "to ensure consistent difficulty across sample sizes",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.5))
add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.35),
            "Configuration", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(6.9), Inches(2.05), Inches(5.7), Inches(1.6), [
    "Share counts tested: 50, 200, 1K, 5K, 20K",
    "Engagements per share: 1, 3, 5, 10",
    "3 random trials per setting",
    "Min 10 engagements per share to be eligible",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(4.2), Inches(6.2), Inches(2.4))
add_textbox(slide, Inches(6.9), Inches(4.3), Inches(5.5), Inches(0.35),
            "Note on Model Strength", size=18, color=RED, bold=True)
add_bullets(slide, Inches(6.9), Inches(4.75), Inches(5.7), Inches(1.6), [
    "Initial run used weak model (50 trees) for speed",
    "Updated to full-strength model (500 trees)",
    "matching the profile predictor config:",
    "     lr=0.05, 127 leaves, max_depth=10",
    "Re-run will give more realistic estimates",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Power Analysis — Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Power Analysis: Results (Initial Run)",
            "50-tree model -- updated 500-tree run pending")
section_label(slide, "Power Analysis - Results")

metric_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(1.0),
            "Job Role: Min Shares", "5,000", RED)
metric_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(1.0),
            "Industry: Min Shares", "20,000", BLUE)
metric_card(slide, Inches(8.7), Inches(1.5), Inches(3.8), Inches(1.0),
            "Company: Min Shares", "200", RGBColor(0x2E, 0xCC, 0x71))

# Summary text below metrics
add_card(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.5))
add_bullets(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.2), [
    "Recommendation: Track at least 20,000 shares, scrape 3-5 profiles per share (~80K total profiles)",
    "These results are from the 50-tree model -- the updated 500-tree run will show higher accuracy",
    "The full profile predictor (500 trees) already shows 0.546 / 0.726 / 0.908 top-3 for these targets",
], size=13, color=DARK)

# Individual power analysis figure slides
pa_slide_figs = [
    ("learning_curve_shares.png", "How Many Shares Do You Need?", "Power Analysis - Shares"),
    ("learning_curve_engagements.png", "How Many Engagements per Share?", "Power Analysis - Engagements"),
    ("heatmap.png", "Top-3 Accuracy: Shares vs Engagements", "Power Analysis - Heatmap"),
]

for fname, title, label in pa_slide_figs:
    fpath = os.path.join(slide_figs_dir, fname)
    if os.path.exists(fpath):
        slide = prs.slides.add_slide(LY_BLANK)
        slide_title(slide, title)
        section_label(slide, label)
        slide.shapes.add_picture(fpath, Inches(1.5), Inches(1.3), Inches(10), Inches(5.8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Model 3 — Method
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 3: Audience Distribution Predictor",
            "Predicting the mix of audience segments per share")
section_label(slide, "Model 3 - Method")

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.35),
            "Key Insight", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.3), Inches(4.2), [
    "Models 1 & 2 predict at the individual level",
    "",
    "But the real business question is:",
    "'What audience mix will this post attract?'",
    "",
    "Distribution prediction is more natural:",
    "     Operates at share level (1 row per share)",
    "     Target is a proportion vector per category",
    "     E.g., '40% Marketing, 30% Sales, 20% Eng'",
    "     Less noisy than individual predictions",
    "     Directly actionable for content strategy",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.5))
add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.35),
            "Method", size=18, color=RGBColor(0x2E, 0xCC, 0x71), bold=True)
add_bullets(slide, Inches(6.9), Inches(2.05), Inches(5.7), Inches(1.6), [
    "Aggregate engagements per share into proportion vectors",
    "Filter shares with >= 5 complete-profile engagements",
    "Multi-output regression (LightGBM per category)",
    "Clip negatives and re-normalize to valid distributions",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(4.2), Inches(6.2), Inches(2.4))
add_textbox(slide, Inches(6.9), Inches(4.3), Inches(5.5), Inches(0.35),
            "Evaluation Metrics", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(6.9), Inches(4.75), Inches(5.7), Inches(1.6), [
    "Cosine Similarity (1.0 = perfect match)",
    "Jensen-Shannon Divergence (0.0 = identical)",
    "Top-K Category Overlap (dominant categories correct?)",
    "Proportion MAE (average error per category)",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Model 3 — Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Model 3: Results", "Audience Distribution Predictor (155K train / 31K test)")
section_label(slide, "Model 3 - Results")

# Metric cards row
metric_card(slide, Inches(0.3), Inches(1.5), Inches(3.0), Inches(1.0),
            "Best Cosine Similarity", "0.796", BLUE)
metric_card(slide, Inches(3.5), Inches(1.5), Inches(3.0), Inches(1.0),
            "Best JS Divergence", "0.438", RGBColor(0x2E, 0xCC, 0x71))
metric_card(slide, Inches(6.7), Inches(1.5), Inches(3.0), Inches(1.0),
            "Best Top-1 Match", "79.2%", YELLOW)
metric_card(slide, Inches(9.9), Inches(1.5), Inches(3.0), Inches(1.0),
            "Shares Used", "194K", RED)

# Results table
headers = ["Dimension", "Cosine Sim", "JS Div", "MAE", "Top-1", "Top-3", "Top-5"]
col_lefts = [Inches(0.8), Inches(2.8), Inches(4.6), Inches(6.1), Inches(7.6), Inches(9.1), Inches(10.6)]
table_top = Inches(2.9)

add_card(slide, Inches(0.5), table_top, Inches(12.3), Inches(0.55), fill=BLUE)
for pos, hdr in zip(col_lefts, headers):
    add_textbox(slide, pos, table_top + Inches(0.08), Inches(1.8), Inches(0.35),
                hdr, size=14, color=WHITE, bold=True)

dist_targets = [
    ("Job Role", "25", "0.757", "0.445", "0.0366", "0.541", "0.557", "0.541", RED),
    ("Industry", "114", "0.770", "0.467", "0.0083", "0.616", "0.557", "0.456", BLUE),
    ("Company", "192", "0.796", "0.438", "0.0047", "0.792", "0.296*", "0.190*", RGBColor(0x2E, 0xCC, 0x71)),
    ("Job Level", "TBD", "TBD", "TBD", "TBD", "TBD", "TBD", "TBD", RGBColor(0x9B, 0x59, 0xB6)),
    ("Country", "TBD", "TBD", "TBD", "TBD", "TBD", "TBD", "TBD", RGBColor(0xF3, 0x9C, 0x12)),
]

for i, (name, ncats, cos, js, mae, t1, t3, t5, accent) in enumerate(dist_targets):
    row_top = table_top + Inches(0.6 + i * 0.6)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, Inches(0.5), row_top, Inches(12.3), Inches(0.55), fill=bg)
    vals = [f"{name} ({ncats})", cos, js, mae, t1, t3, t5]
    for pos, val in zip(col_lefts, vals):
        c = accent if pos == col_lefts[0] else DARK
        b = pos == col_lefts[0]
        add_textbox(slide, pos, row_top + Inches(0.08), Inches(1.8), Inches(0.35),
                    val, size=14, color=c, bold=b)

# Insights
add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8))
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.35),
            "Analysis", size=18, color=DARK, bold=True)
add_bullets(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.2), [
    "Company distribution has the highest cosine similarity (0.796) and best top-1 match (79.2%)",
    "Proportion MAE is very low across all targets -- predicted distributions closely match actual",
    "*Regressor top-K overlap drops for company because 192 near-zero proportions cause ranking noise",
    "Classifier approach fixes this: 72.6% top-3 for company at 20K shares (see sample analysis)",
], size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES: Model 3 — Distribution Figures
# ═══════════════════════════════════════════════════════════════════════════════
dist_fig_dir = "./figures/distribution"
dist_figs = [
    ("dist_comparison_job_title.png", "Predicted vs Actual: Job Title Distribution", "Model 3 - Job Title"),
    ("dist_comparison_industry.png", "Predicted vs Actual: Industry Distribution", "Model 3 - Industry"),
    ("dist_comparison_company.png", "Predicted vs Actual: Company Distribution", "Model 3 - Company"),
    ("dist_comparison_job_level.png", "Predicted vs Actual: Job Level Distribution", "Model 3 - Job Level"),
    ("dist_comparison_country.png", "Predicted vs Actual: Country Distribution", "Model 3 - Country"),
    ("scatter_top1_proportion.png", "Top-1 Category: Predicted vs Actual Proportion", "Model 3 - Scatter"),
    ("metrics_summary.png", "Distribution Prediction: Metrics Summary", "Model 3 - Summary"),
]

for fname, title, label in dist_figs:
    fpath = os.path.join(dist_fig_dir, fname)
    if os.path.exists(fpath):
        slide = prs.slides.add_slide(LY_BLANK)
        slide_title(slide, title)
        section_label(slide, label)
        slide.shapes.add_picture(fpath, Inches(1.5), Inches(1.3), Inches(10), Inches(5.8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE: Distribution Sample Size & Model Comparison — Method
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Distribution Predictor: Sample Size & Model Comparison",
            "Improving top-K accuracy with a classifier approach")
section_label(slide, "Model 3 - Sample Analysis")

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.35),
            "Why Classifier?", size=18, color=RED, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.3), Inches(4.2), [
    "Regressor: MultiOutputRegressor trains one",
    "  model per category independently",
    "  Small proportion errors shuffle rankings",
    "  -> poor top-K overlap",
    "",
    "Classifier: LGBMClassifier with softmax",
    "  predicts 'which category dominates?'",
    "  Probabilities are naturally ranked",
    "  -> directly optimises for top-K",
    "",
    "Same LightGBM algorithm, different framing",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.5))
add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.35),
            "Experiment Design", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(6.9), Inches(2.05), Inches(5.7), Inches(1.6), [
    "Sample sizes: 5,000 / 10,000 / 20,000 shares",
    "3 trials per setting for confidence intervals",
    "Both regressor and classifier on same splits",
    "Time-based 80/20 train/test split",
], size=13, color=DARK)

add_card(slide, Inches(6.6), Inches(4.2), Inches(6.2), Inches(2.4))
add_textbox(slide, Inches(6.9), Inches(4.3), Inches(5.5), Inches(0.35),
            "Metrics Compared", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(6.9), Inches(4.75), Inches(5.7), Inches(1.6), [
    "Top-3 Accuracy: true dominant in top-3 predictions?",
    "Top-5 Accuracy: true dominant in top-5 predictions?",
    "Evaluated at each sample size for both approaches",
    "LightGBM: 500 trees, lr=0.05, 127 leaves",
], size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE: Distribution Sample Size — Top-3 Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Distribution Predictor: Top-3 & Top-5 Accuracy",
            "Regressor vs Classifier across sample sizes")
section_label(slide, "Model 3 - Sample Results")

# Top-3 table
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.35),
            "Top-3 Accuracy", size=18, color=BLUE, bold=True)

headers_3 = ["Dimension", "Approach", "5K", "10K", "20K"]
col_lefts_3 = [Inches(0.8), Inches(3.0), Inches(5.2), Inches(7.0), Inches(8.8)]
table_top_3 = Inches(1.95)

add_card(slide, Inches(0.5), table_top_3, Inches(10.5), Inches(0.45), fill=BLUE)
for pos, hdr in zip(col_lefts_3, headers_3):
    add_textbox(slide, pos, table_top_3 + Inches(0.05), Inches(1.8), Inches(0.3),
                hdr, size=13, color=WHITE, bold=True)

top3_data = [
    ("Job Role", "Regressor", "0.482", "0.509", "0.526", RED),
    ("Job Role", "Classifier", "0.382", "0.411", "0.443", RED),
    ("Industry", "Regressor", "0.441", "0.496", "0.524", BLUE),
    ("Industry", "Classifier", "0.439", "0.530", "0.560", BLUE),
    ("Company", "Regressor", "0.197", "0.244", "0.269", RGBColor(0x2E, 0xCC, 0x71)),
    ("Company", "Classifier", "0.495", "0.658", "0.726", RGBColor(0x2E, 0xCC, 0x71)),
    ("Job Level", "Regressor", "TBD", "TBD", "TBD", RGBColor(0x9B, 0x59, 0xB6)),
    ("Job Level", "Classifier", "TBD", "TBD", "TBD", RGBColor(0x9B, 0x59, 0xB6)),
    ("Country", "Regressor", "TBD", "TBD", "TBD", RGBColor(0xF3, 0x9C, 0x12)),
    ("Country", "Classifier", "TBD", "TBD", "TBD", RGBColor(0xF3, 0x9C, 0x12)),
]

for i, (name, approach, v5k, v10k, v20k, accent) in enumerate(top3_data):
    row_top = table_top_3 + Inches(0.5 + i * 0.4)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, Inches(0.5), row_top, Inches(10.5), Inches(0.38), fill=bg)
    vals = [name, approach, v5k, v10k, v20k]
    for pos, val in zip(col_lefts_3, vals):
        c = accent if pos == col_lefts_3[0] else DARK
        b = pos == col_lefts_3[0] or pos == col_lefts_3[1]
        add_textbox(slide, pos, row_top + Inches(0.03), Inches(1.8), Inches(0.3),
                    val, size=12, color=c, bold=b)

# Top-5 table
add_textbox(slide, Inches(0.8), Inches(4.55), Inches(5), Inches(0.35),
            "Top-5 Accuracy", size=18, color=RGBColor(0x2E, 0xCC, 0x71), bold=True)

table_top_5 = Inches(4.95)
add_card(slide, Inches(0.5), table_top_5, Inches(10.5), Inches(0.45), fill=RGBColor(0x2E, 0xCC, 0x71))
for pos, hdr in zip(col_lefts_3, headers_3):
    add_textbox(slide, pos, table_top_5 + Inches(0.05), Inches(1.8), Inches(0.3),
                hdr, size=13, color=WHITE, bold=True)

top5_data = [
    ("Job Role", "Regressor", "0.487", "0.505", "0.521", RED),
    ("Job Role", "Classifier", "0.382", "0.411", "0.443", RED),
    ("Industry", "Regressor", "0.352", "0.398", "0.424", BLUE),
    ("Industry", "Classifier", "0.439", "0.530", "0.560", BLUE),
    ("Company", "Regressor", "0.131", "0.159", "0.175", RGBColor(0x2E, 0xCC, 0x71)),
    ("Company", "Classifier", "0.495", "0.658", "0.726", RGBColor(0x2E, 0xCC, 0x71)),
    ("Job Level", "Regressor", "TBD", "TBD", "TBD", RGBColor(0x9B, 0x59, 0xB6)),
    ("Job Level", "Classifier", "TBD", "TBD", "TBD", RGBColor(0x9B, 0x59, 0xB6)),
    ("Country", "Regressor", "TBD", "TBD", "TBD", RGBColor(0xF3, 0x9C, 0x12)),
    ("Country", "Classifier", "TBD", "TBD", "TBD", RGBColor(0xF3, 0x9C, 0x12)),
]

for i, (name, approach, v5k, v10k, v20k, accent) in enumerate(top5_data):
    row_top = table_top_5 + Inches(0.5 + i * 0.4)
    bg = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, Inches(0.5), row_top, Inches(10.5), Inches(0.38), fill=bg)
    vals = [name, approach, v5k, v10k, v20k]
    for pos, val in zip(col_lefts_3, vals):
        c = accent if pos == col_lefts_3[0] else DARK
        b = pos == col_lefts_3[0] or pos == col_lefts_3[1]
        add_textbox(slide, pos, row_top + Inches(0.03), Inches(1.8), Inches(0.3),
                    val, size=12, color=c, bold=b)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES: Distribution Sample Analysis — Figures
# ═══════════════════════════════════════════════════════════════════════════════
dist_sa_fig_dir = "./figures/distribution_sample_analysis"
dist_sa_figs = [
    ("learning_curve_top3.png", "Top-3 Accuracy: Regressor vs Classifier by Sample Size", "Model 3 - Top-3 Curve"),
    ("learning_curve_top5.png", "Top-5 Accuracy: Regressor vs Classifier by Sample Size", "Model 3 - Top-5 Curve"),
    ("comparison_bar.png", "Regressor vs Classifier: Side-by-Side Comparison", "Model 3 - Comparison"),
]

for fname, title, label in dist_sa_figs:
    fpath = os.path.join(dist_sa_fig_dir, fname)
    if os.path.exists(fpath):
        slide = prs.slides.add_slide(LY_BLANK)
        slide_title(slide, title)
        section_label(slide, label)
        slide.shapes.add_picture(fpath, Inches(1.5), Inches(1.3), Inches(10), Inches(5.8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Summary & Next Steps
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Summary & Next Steps")
section_label(slide, "Summary")

# Left: What We've Built
add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.35),
            "What We've Built", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.5), Inches(1.8), [
    "Engagement volume predictor -- 95% recall (ROC-AUC 0.95)",
    "Engager profile predictor -- 83% top-1 company, 73% top-3 industry",
    "Audience distribution predictor -- 0.796 cosine sim, 79% top-1 (company)",
    "Power analysis -- quantifies scraping data requirements",
], size=13, color=DARK)

# Right: Next Steps
add_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5))
add_textbox(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.35),
            "Next Steps", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(7.1), Inches(2.05), Inches(5.5), Inches(1.8), [
    "Run profile predictor on full dataset (currently 50K sample)",
    "Run updated power analysis with 500-tree model",
    "Compare distribution vs classification approach",
    "Explore sentence transformer embeddings for text",
    "Productionize: API endpoint for real-time predictions",
], size=13, color=DARK)

# Bottom: Key Takeaway
add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.3))
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.35),
            "Key Takeaway", size=18, color=RGBColor(0x2E, 0xCC, 0x71), bold=True)
add_bullets(slide, Inches(0.8), Inches(4.85), Inches(11.5), Inches(1.5), [
    "We can predict both HOW MUCH and WHO engages with LinkedIn shares using share content and sharer profile",
    "Company/employer is easiest to predict (strong client signal), industry is moderate, job title is hardest",
    "The distribution approach (Model 3) directly answers 'what audience mix?' with 0.76-0.80 cosine similarity",
    "If switching to scraped data: ~20K shares with 3-5 profiles each provides near-maximum predictive power",
], size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE: Why Job Role Is Hard — And How To Improve It
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(LY_BLANK)
slide_title(slide, "Why Job Role Is Hard to Predict",
            "Root causes and actionable improvements")
section_label(slide, "Job Role Analysis")

# Left card: Why it's hard
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1))
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.3), Inches(0.35),
            "Root Causes", size=18, color=RED, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.05), Inches(5.3), Inches(4.2), [
    "Job roles cross-cut industries and topics",
    "  A 'leadership' post attracts VPs from every",
    "  industry -- no clean content signal",
    "  Unlike company (network-clustered) or industry",
    "  (topic-correlated), role is not content-predictable",
    "",
    "TF-IDF misses role-specific semantics",
    "  'Pipeline generation' and 'quota attainment'",
    "  both signal Sales -- TF-IDF treats them as",
    "  unrelated, sentence embeddings would not",
    "",
    "Noisy ground truth at share level",
    "  With only 5+ complete-profile engagements,",
    "  a single random engagement shifts proportions",
    "  by 20% -- too noisy for role prediction",
], size=12, color=DARK)

# Right top: High-impact improvements
add_card(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(2.8))
add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.35),
            "High Impact", size=18, color=BLUE, bold=True)
add_bullets(slide, Inches(6.9), Inches(2.05), Inches(5.8), Inches(2.1), [
    "Sentence transformer embeddings",
    "  Replace TF-IDF with all-MiniLM-L6-v2",
    "  Captures semantic role language end-to-end",
    "",
    "Increase min engagements per share (5 -> 15+)",
    "  Cleaner ground truth distributions",
    "  Fewer shares but much less label noise",
], size=12, color=DARK)

# Right bottom: Moderate-impact improvements
add_card(slide, Inches(6.6), Inches(4.5), Inches(6.2), Inches(2.1))
add_textbox(slide, Inches(6.9), Inches(4.6), Inches(5.8), Inches(0.35),
            "Moderate Impact", size=18, color=YELLOW, bold=True)
add_bullets(slide, Inches(6.9), Inches(5.05), Inches(5.8), Inches(1.4), [
    "Sharer role x content type interaction features",
    "  A Sales VP sharing content attracts a different",
    "  role mix than an Engineering VP sharing the same",
    "Hierarchical grouping (25 -> 8 broad buckets)",
    "  Technical / Commercial / Marketing / Executive...",
], size=12, color=DARK)

# ─── Save ────────────────────────────────────────────────────────────────────
out_path = "./engagement_prediction_deck.pptx"
prs.save(out_path)
print(f"Saved presentation to: {out_path}")
